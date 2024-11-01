from pptx import Presentation
from pptx.util import Inches
from datetime import datetime

# Maak een nieuwe presentatie aan
prs = Presentation('test1.pptx')
img_path = 'Spie-vector-logo.png'
now = datetime.now()

# Voeg een dia toe voor de titel
slide_layout = prs.slide_layouts[0]  # 0 is de layout voor een titelpagina
slide = prs.slides[0]

# Voeg titel en ondertitel toe aan de titelpagina
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = f"Maandrapportage van: {now:%B %Y}"
subtitle.text = "Gegenereerd met Python"

left = Inches(8)
top = Inches(.2)
height = Inches(1)
pic = slide.shapes.add_picture(img_path, left, top, height=height)

# Voeg een nieuwe dia toe voor inhoud
slide_layout = prs.slide_layouts[1]  # 1 is de layout voor een titel + inhoud
slide = prs.slides[1]

# Voeg een titel en inhoud toe aan de inhoudsdia
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Overzicht"
content.text = "Dit is een eenvoudige dia met wat tekst."

# Sla de presentatie op
prs.save('test1.pptx')
